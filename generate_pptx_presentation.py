import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to 16:9 widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    # Theme Palette (Modern Dark Academic Theme)
    BG_COLOR = RGBColor(15, 23, 42)       # Slate Navy #0F172A
    CARD_BG = RGBColor(30, 41, 59)        # Dark Card Slate #1E293B
    BORDER_COLOR = RGBColor(51, 65, 85)   # Border #334155
    TEXT_MAIN = RGBColor(248, 250, 252)   # Pure White #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate Muted #94A3B8
    CYAN_ACCENT = RGBColor(56, 189, 248)  # Bright Cyan #38BDF8
    EMERALD_ACCENT = RGBColor(52, 211, 153) # Emerald Green #34D399
    AMBER_ACCENT = RGBColor(251, 191, 36)  # Warm Amber #FBBF24
    VIOLET_ACCENT = RGBColor(168, 85, 247) # Electric Violet #A855F7
    
    fig_dir = r"d:\college\DL 4 models\zz paper\figures"

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="RESEARCH PAPER PRESENTATION"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN_ACCENT
        p_cat.font.name = "Calibri"
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN
        p_title.font.name = "Calibri"

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_COLOR):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # ==========================================
    # SLIDE 1: Title Slide (Hero Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_bg(slide1)
    
    add_card(slide1, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1), CARD_BG, CYAN_ACCENT)
    
    tb = slide1.shapes.add_textbox(Inches(1.3), Inches(1.5), Inches(10.7), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "RESEARCH PAPER PRESENTATION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = "A Comparison of Spatial Attention, Multi-Scale Fusion, and Deep Residual Models on Facial Emotion Recognition"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN
    p2.space_before = Pt(12)
    p2.space_after = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "All 7 Result Values (Acc, Prec, Rec, Spec, F1, Kappa, AUC) Across All Stages: Base → Fine-Tune → Hyperparameter → XAI & Ablation"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_after = Pt(20)
    
    p4 = tf.add_paragraph()
    p4.text = "Author: Vijay Mahesh  |  Affiliation: Department of Deep Learning, Anna University"
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = EMERALD_ACCENT

    # ==========================================
    # SLIDE 2: Executive Summary & Stage Highlights
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide2)
    add_header(slide2, "Executive Summary & All-Stage Metric Evolution", "Overview")
    
    stats = [
        ("98.57%", "CHAMPION ACCURACY & F1 (0.9857)", "Multi-Scale Fused EfficientNet-B0 final state across 7 emotion classes", CYAN_ACCENT),
        ("0.9976", "SPECIFICITY (0.9976)", "Near-flawless True Negative recognition avoiding false positive errors", EMERALD_ACCENT),
        ("0.9984", "AUC-ROC (0.9984) & KAPPA (0.9833)", "Near-perfect class separability and inter-rater agreement over random chance", AMBER_ACCENT)
    ]
    
    for i, (val, title, desc, col) in enumerate(stats):
        left = Inches(0.8 + i * 3.9)
        add_card(slide2, left, Inches(1.6), Inches(3.7), Inches(2.2), CARD_BG, col)
        
        tb = slide2.shapes.add_textbox(left + Inches(0.2), Inches(1.75), Inches(3.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(4)

    add_card(slide2, Inches(0.8), Inches(4.1), Inches(11.733), Inches(2.7), CARD_BG, BORDER_COLOR)
    tb_summary = slide2.shapes.add_textbox(Inches(1.1), Inches(4.3), Inches(11.1), Inches(2.3))
    tf_sum = tb_summary.text_frame
    tf_sum.word_wrap = True
    
    bullets = [
        "Stage 1 (Base Model): Frozen backbones yielded baseline performance (Accuracy: 32.95% - 41.92%, F1: 0.330 - 0.420, AUC: 0.684 - 0.728).",
        "Stage 2 (Fine-Tuning): Unfreezing Conv blocks surged performance (Accuracy: 78.25% - 82.01%, F1: 0.783 - 0.821, AUC: 0.932 - 0.948).",
        "Stage 3 & 4 (Hyperparameter & Champion State): Multi-scale fusion (f_fused ∈ ℝ²¹⁹²) reached 98.57% Accuracy, 0.9859 Precision, 0.9857 Recall, 0.9976 Specificity, 0.9857 F1, 0.9833 Kappa, & 0.9984 AUC-ROC."
    ]
    for idx, b in enumerate(bullets):
        p = tf_sum.add_paragraph() if idx > 0 else tf_sum.paragraphs[0]
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: Introduction & Biological Basis
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_bg(slide3)
    add_header(slide3, "Introduction: Affective Computing & Biological FACS Basis", "Introduction")
    
    add_card(slide3, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2), CARD_BG, CYAN_ACCENT)
    tb_left = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "🧬 Facial Action Coding System (FACS)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.space_after = Pt(12)
    
    l_points = [
        "Facial Emotion Recognition (FER) automates emotional state classification from visual streams.",
        "Anatomical Grounding: Facial expressions stem from contractions of key underlying facial muscles:",
        "  - Zygomaticus Major: Controls mouth corner pull (Happy / Smiling).",
        "  - Corrugator Supercilii: Controls eyebrow furrowing (Angry / Fear).",
        "FACS maps these physical Action Units (AUs) into numerical spatial coordinates.",
        "Core Goal: Simultaneous capture of fine local muscular deformations and deep global facial semantics."
    ]
    for pt in l_points:
        p = tf_l.add_paragraph()
        p.text = pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(6)

    add_card(slide3, Inches(6.833), Inches(1.6), Inches(5.7), Inches(5.2), CARD_BG, EMERALD_ACCENT)
    tb_right = slide3.shapes.add_textbox(Inches(7.033), Inches(1.8), Inches(5.3), Inches(4.8))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "🚀 Key Real-World Application Domains"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT
    p.space_after = Pt(12)
    
    r_points = [
        "Human-Computer Interaction (HCI): Adaptive e-learning platforms and empathetic virtual tutors.",
        "Automotive Safety: Non-intrusive driver fatigue, drowsiness, and road-rage stress monitoring.",
        "Healthcare & Psychiatry: Automated diagnostic support for autism spectrum and depression tracking.",
        "Consumer Market Research: Real-time audience sentiment tracking during video previews.",
        "System Requirement: Robust performance across illumination changes and head tilts."
    ]
    for pt in r_points:
        p = tf_r.add_paragraph()
        p.text = pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 4: Problem Statement
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_bg(slide4)
    add_header(slide4, "Problem Statement & Critical Technical Challenges", "Problem Statement")
    
    probs = [
        ("01", "Spatial Detail Degradation in Deep CNNs", "Standard CNNs rely on consecutive max-pooling operations. As spatial resolution drops from 224x224 down to 7x7, local spatial micro-expression contours (eyebrow curvature, lip displacement) are erased before final classification.", AMBER_ACCENT),
        ("02", "Inadequacy of Traditional Hand-Crafted Features", "Hand-crafted descriptors (LBP, Gabor filters + SVR/XGBoost) are computationally lightweight but brittle. They fail to generalize under illumination fluctuations, head tilts, and non-linear facial poses.", CYAN_ACCENT),
        ("03", "Accuracy vs. Real-Time Edge Latency Trade-off", "Heavy deep learning models achieve high accuracy but demand immense compute resources. Lightweight models fit edge nodes but sacrifice accuracy. A multi-level architectural remedy is required.", VIOLET_ACCENT)
    ]
    
    for i, (num, title, desc, col) in enumerate(probs):
        top = Inches(1.6 + i * 1.8)
        add_card(slide4, Inches(0.8), top, Inches(11.733), Inches(1.6), CARD_BG, col)
        
        add_card(slide4, Inches(1.0), top + Inches(0.2), Inches(0.8), Inches(1.2), col, None)
        tb_num = slide4.shapes.add_textbox(Inches(1.0), top + Inches(0.4), Inches(0.8), Inches(0.8))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = num
        p_num.font.size = Pt(20)
        p_num.font.bold = True
        p_num.font.color.rgb = BG_COLOR
        p_num.alignment = PP_ALIGN.CENTER
        
        tb_desc = slide4.shapes.add_textbox(Inches(2.0), top + Inches(0.15), Inches(10.3), Inches(1.3))
        tf_d = tb_desc.text_frame
        tf_d.word_wrap = True
        
        p = tf_d.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf_d.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 5: Level 0 - Balanced Dataset & Deep Analysis
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_bg(slide5)
    add_header(slide5, "Level 0: Balanced Dataset & Preprocessing Pipeline", "Level 0: Dataset")
    
    add_card(slide5, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2), CARD_BG, EMERALD_ACCENT)
    tb_ds = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf_ds = tb_ds.text_frame
    tf_ds.word_wrap = True
    
    p = tf_ds.paragraphs[0]
    p.text = "⚖️ Dataset Balancing & Preprocessing"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT
    p.space_after = Pt(10)
    
    ds_pts = [
        "Dataset Volume: 7,529 facial images across 7 emotion classes.",
        "Initial Imbalance:",
        "  • Disgust (Minority): 460 images vs Happy (Majority): 1,197 images",
        "  • Angry: 1,186 | Fear: 1,188 | Neutral: 1,194 | Sad: 1,189 | Surprise: 1,115",
        "Contrast Limited Adaptive Histogram Equalization (CLAHE):",
        "  • Clip limit: 2.0 | Grid size: 8x8 — resolves local lighting variance.",
        "Recursive Class Weighting Mechanism:",
        "  • w_c = N / (C * n_c) — scales loss gradients to prevent majority bias.",
        "Stratified Split: 80% Training / 20% Reserved Test (1,050 balanced samples)."
    ]
    for pt in ds_pts:
        p = tf_ds.add_paragraph()
        p.text = pt
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(4)

    add_card(slide5, Inches(6.833), Inches(1.6), Inches(5.7), Inches(5.2), CARD_BG, BORDER_COLOR)
    fig_ds = os.path.join(fig_dir, "fig1_dataset.png")
    if os.path.exists(fig_ds):
        slide5.shapes.add_picture(fig_ds, Inches(7.033), Inches(1.8), Inches(5.3), Inches(4.8))

    # ==========================================
    # SLIDE 6: ALL 7 RESULT METRICS FOR STAGE 1 (BASE MODEL)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_bg(slide6)
    add_header(slide6, "Stage 1: Base Model Phase (All 7 Result Values)", "Stage 1: Base Model")
    
    stage1_models = [
        ("EfficientNet-B0 (Base)", [
            ("Accuracy", "41.59%"), ("Precision", "0.4180"), ("Recall", "0.4159"),
            ("Specificity", "0.9026"), ("F1-Score", "0.4165"), ("Kappa Score", "0.3185"), ("AUC-ROC", "0.7250")
        ], CYAN_ACCENT),
        ("YOLOv8 Classifier (Base)", [
            ("Accuracy", "32.95%"), ("Precision", "0.3312"), ("Recall", "0.3295"),
            ("Specificity", "0.8882"), ("F1-Score", "0.3301"), ("Kappa Score", "0.2178"), ("AUC-ROC", "0.6840")
        ], EMERALD_ACCENT),
        ("ResNet-50 (Base)", [
            ("Accuracy", "41.92%"), ("Precision", "0.4215"), ("Recall", "0.4192"),
            ("Specificity", "0.9032"), ("F1-Score", "0.4200"), ("Kappa Score", "0.3224"), ("AUC-ROC", "0.7285")
        ], AMBER_ACCENT),
        ("MobileNetV2 (Base)", [
            ("Accuracy", "37.27%"), ("Precision", "0.3750"), ("Recall", "0.3727"),
            ("Specificity", "0.8955"), ("F1-Score", "0.3735"), ("Kappa Score", "0.2682"), ("AUC-ROC", "0.7020")
        ], VIOLET_ACCENT)
    ]
    
    for i, (mname, mlist, col) in enumerate(stage1_models):
        left = Inches(0.8 + i * 2.95)
        add_card(slide6, left, Inches(1.5), Inches(2.8), Inches(5.4), CARD_BG, col)
        
        tb = slide6.shapes.add_textbox(left + Inches(0.15), Inches(1.65), Inches(2.5), Inches(5.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = mname
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = "STAGE 1: FROZEN BACKBONE"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(2)
        p2.space_after = Pt(8)
        
        for metric_name, val in mlist:
            p_m = tf.add_paragraph()
            p_m.text = f"• {metric_name}: "
            p_m.font.size = Pt(11)
            p_m.font.color.rgb = TEXT_MUTED
            
            run = p_m.add_run()
            run.text = val
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = col if metric_name in ["Accuracy", "F1-Score", "AUC-ROC"] else TEXT_MAIN
            p_m.space_before = Pt(2)

    # ==========================================
    # SLIDE 7: ALL 7 RESULT METRICS FOR STAGE 2 (FINE-TUNED MODEL)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_bg(slide7)
    add_header(slide7, "Stage 2: Fine-Tuned Model Phase (All 7 Result Values)", "Stage 2: Fine-Tuned")
    
    stage2_models = [
        ("EfficientNet-B0 (Fine-Tuned)", [
            ("Accuracy", "78.25%"), ("Precision", "0.7850"), ("Recall", "0.7825"),
            ("Specificity", "0.9637"), ("F1-Score", "0.7832"), ("Kappa Score", "0.7462"), ("AUC-ROC", "0.9320")
        ], CYAN_ACCENT),
        ("YOLOv8 Classifier (Fine-Tuned)", [
            ("Accuracy", "82.01%"), ("Precision", "0.8220"), ("Recall", "0.8201"),
            ("Specificity", "0.9700"), ("F1-Score", "0.8208"), ("Kappa Score", "0.7901"), ("AUC-ROC", "0.9480")
        ], EMERALD_ACCENT),
        ("ResNet-50 (Fine-Tuned)", [
            ("Accuracy", "79.92%"), ("Precision", "0.8015"), ("Recall", "0.7992"),
            ("Specificity", "0.9665"), ("F1-Score", "0.8001"), ("Kappa Score", "0.7657"), ("AUC-ROC", "0.9395")
        ], AMBER_ACCENT),
        ("MobileNetV2 (Fine-Tuned)", [
            ("Accuracy", "80.12%"), ("Precision", "0.8035"), ("Recall", "0.8012"),
            ("Specificity", "0.9668"), ("F1-Score", "0.8020"), ("Kappa Score", "0.7680"), ("AUC-ROC", "0.9410")
        ], VIOLET_ACCENT)
    ]
    
    for i, (mname, mlist, col) in enumerate(stage2_models):
        left = Inches(0.8 + i * 2.95)
        add_card(slide7, left, Inches(1.5), Inches(2.8), Inches(5.4), CARD_BG, col)
        
        tb = slide7.shapes.add_textbox(left + Inches(0.15), Inches(1.65), Inches(2.5), Inches(5.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = mname
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = "STAGE 2: CONV UNFREEZING"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(2)
        p2.space_after = Pt(8)
        
        for metric_name, val in mlist:
            p_m = tf.add_paragraph()
            p_m.text = f"• {metric_name}: "
            p_m.font.size = Pt(11)
            p_m.font.color.rgb = TEXT_MUTED
            
            run = p_m.add_run()
            run.text = val
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = col if metric_name in ["Accuracy", "F1-Score", "AUC-ROC"] else TEXT_MAIN
            p_m.space_before = Pt(2)

    # ==========================================
    # SLIDE 8: ALL 7 RESULT METRICS FOR STAGE 3 (HYPERPARAMETER MODEL)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_bg(slide8)
    add_header(slide8, "Stage 3: Hyperparameter Model Phase (All 7 Result Values)", "Stage 3: Hyperparameter")
    
    stage3_models = [
        ("EfficientNet-B0 (Grid Search)", [
            ("Accuracy", "89.41%"), ("Precision", "0.8965"), ("Recall", "0.8941"),
            ("Specificity", "0.9823"), ("F1-Score", "0.8950"), ("Kappa Score", "0.8764"), ("AUC-ROC", "0.9725")
        ], CYAN_ACCENT),
        ("YOLOv8 Classifier (Grid Search)", [
            ("Accuracy", "94.26%"), ("Precision", "0.9438"), ("Recall", "0.9426"),
            ("Specificity", "0.9904"), ("F1-Score", "0.9430"), ("Kappa Score", "0.9330"), ("AUC-ROC", "0.9860")
        ], EMERALD_ACCENT),
        ("ResNet-50 (Grid Search)", [
            ("Accuracy", "87.95%"), ("Precision", "0.8812"), ("Recall", "0.8795"),
            ("Specificity", "0.9799"), ("F1-Score", "0.8801"), ("Kappa Score", "0.8594"), ("AUC-ROC", "0.9680")
        ], AMBER_ACCENT),
        ("MobileNetV2 (Grid Search)", [
            ("Accuracy", "87.34%"), ("Precision", "0.8755"), ("Recall", "0.8734"),
            ("Specificity", "0.9789"), ("F1-Score", "0.8742"), ("Kappa Score", "0.8523"), ("AUC-ROC", "0.9650")
        ], VIOLET_ACCENT)
    ]
    
    for i, (mname, mlist, col) in enumerate(stage3_models):
        left = Inches(0.8 + i * 2.95)
        add_card(slide8, left, Inches(1.5), Inches(2.8), Inches(5.4), CARD_BG, col)
        
        tb = slide8.shapes.add_textbox(left + Inches(0.15), Inches(1.65), Inches(2.5), Inches(5.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = mname
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = "STAGE 3: GRID SEARCH TUNED"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(2)
        p2.space_after = Pt(8)
        
        for metric_name, val in mlist:
            p_m = tf.add_paragraph()
            p_m.text = f"• {metric_name}: "
            p_m.font.size = Pt(11)
            p_m.font.color.rgb = TEXT_MUTED
            
            run = p_m.add_run()
            run.text = val
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = col if metric_name in ["Accuracy", "F1-Score", "AUC-ROC"] else TEXT_MAIN
            p_m.space_before = Pt(2)

    # ==========================================
    # SLIDE 9: ALL 7 RESULT METRICS FOR STAGE 4 (CHAMPION & XAI/ABLATION STATE)
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_bg(slide9)
    add_header(slide9, "Stage 4: Champion & XAI/Ablation State (All 7 Result Values)", "Stage 4: Champion State")
    
    stage4_models = [
        ("EfficientNet-B0 🏆", [
            ("Accuracy", "98.57%"), ("Precision", "0.9859"), ("Recall", "0.9857"),
            ("Specificity", "0.9976"), ("F1-Score", "0.9857"), ("Kappa Score", "0.9833"), ("AUC-ROC", "0.9984")
        ], CYAN_ACCENT),
        ("YOLOv8 Classifier", [
            ("Accuracy", "95.52%"), ("Precision", "0.9553"), ("Recall", "0.9552"),
            ("Specificity", "0.9925"), ("F1-Score", "0.9552"), ("Kappa Score", "0.9478"), ("AUC-ROC", "0.9900")
        ], EMERALD_ACCENT),
        ("ResNet-50", [
            ("Accuracy", "94.57%"), ("Precision", "0.9462"), ("Recall", "0.9457"),
            ("Specificity", "0.9910"), ("F1-Score", "0.9458"), ("Kappa Score", "0.9367"), ("AUC-ROC", "0.9825")
        ], AMBER_ACCENT),
        ("MobileNetV2", [
            ("Accuracy", "93.52%"), ("Precision", "0.9359"), ("Recall", "0.9352"),
            ("Specificity", "0.9892"), ("F1-Score", "0.9353"), ("Kappa Score", "0.9244"), ("AUC-ROC", "0.9871")
        ], VIOLET_ACCENT)
    ]
    
    for i, (mname, mlist, col) in enumerate(stage4_models):
        left = Inches(0.8 + i * 2.95)
        add_card(slide9, left, Inches(1.5), Inches(2.8), Inches(5.4), CARD_BG, col)
        
        tb = slide9.shapes.add_textbox(left + Inches(0.15), Inches(1.65), Inches(2.5), Inches(5.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = mname
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = "STAGE 4: CHAMPION CONVERGED"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(2)
        p2.space_after = Pt(8)
        
        for metric_name, val in mlist:
            p_m = tf.add_paragraph()
            p_m.text = f"• {metric_name}: "
            p_m.font.size = Pt(11)
            p_m.font.color.rgb = TEXT_MUTED
            
            run = p_m.add_run()
            run.text = val
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = col if metric_name in ["Accuracy", "F1-Score", "AUC-ROC"] else TEXT_MAIN
            p_m.space_before = Pt(2)

    # ==========================================
    # SLIDE 10: STAGE-BY-STAGE PROGRESSION TABLE FOR CHAMPION EFFICIENTNET-B0
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_bg(slide10)
    add_header(slide10, "Champion EfficientNet-B0: Stage-by-Stage Metric Evolution", "Champion Trajectory")
    
    rows, cols = 8, 5
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3)
    table_shape = slide10.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Evaluation Metric", "Stage 1: Base Model", "Stage 2: Fine-Tuned", "Stage 3: Grid Search", "Stage 4: Champion State"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CYAN_ACCENT if j == 4 else CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = BG_COLOR if j == 4 else TEXT_MAIN
        p.alignment = PP_ALIGN.CENTER
        
    champ_stage_data = [
        ["Accuracy", "41.59%", "78.25%", "89.41%", "98.57% 🏆"],
        ["Precision (Macro)", "0.4180", "0.7850", "0.8965", "0.9859"],
        ["Recall (Macro)", "0.4159", "0.7825", "0.8941", "0.9857"],
        ["Specificity (Macro)", "0.9026", "0.9637", "0.9823", "0.9976"],
        ["F1-Score (Macro)", "0.4165", "0.7832", "0.8950", "0.9857"],
        ["Cohen's Kappa (κ)", "0.3185", "0.7462", "0.8764", "0.9833"],
        ["AUC-ROC (OVR)", "0.7250", "0.9320", "0.9725", "0.9984"]
    ]
    
    for i, row in enumerate(champ_stage_data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if j == 4:
                p.font.bold = True
                p.font.color.rgb = CYAN_ACCENT
            else:
                p.font.color.rgb = TEXT_MAIN
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

    # ==========================================
    # SLIDE 11: Level 4A - XAI Analysis
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    apply_bg(slide11)
    add_header(slide11, "Level 4A: Explainable AI (XAI) Multi-Layer Verification", "Level 4A: XAI")
    
    add_card(slide11, Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.2), CARD_BG, CYAN_ACCENT)
    tb_xai = slide11.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.8), Inches(4.8))
    tf_x = tb_xai.text_frame
    tf_x.word_wrap = True
    
    p = tf_x.paragraphs[0]
    p.text = "🔍 Multi-Layer XAI Techniques"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.space_after = Pt(10)
    
    xai_pts = [
        "1. Grad-CAM (Gradient Class Activation):",
        "  • Target Layer: top_conv / Conv_1 / conv5_block3_out.",
        "  • Highlights high activation strictly on eyes, eyebrows, and mouth contour.",
        "2. Occlusion Sensitivity Analysis:",
        "  • Systematic 32x32 pixel patch occlusion.",
        "  • Confirms maximal confidence drop when eyes/mouth are covered.",
        "3. YOLOv8 9-Stage Feature Maps:",
        "  • Hierarchical compression: Stage 0 (Edges 548KB) → Stage 8 (Semantics 20KB)."
    ]
    for pt in xai_pts:
        p = tf_x.add_paragraph()
        p.text = pt
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(6)

    add_card(slide11, Inches(6.3), Inches(1.6), Inches(6.233), Inches(5.2), CARD_BG, BORDER_COLOR)
    fig_gcam = os.path.join(fig_dir, "figure13_gradcam.png")
    if os.path.exists(fig_gcam):
        slide11.shapes.add_picture(fig_gcam, Inches(6.5), Inches(1.8), Inches(5.833), Inches(4.8))

    # ==========================================
    # SLIDE 12: Level 4B - Ablation Study Level-by-Level
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    apply_bg(slide12)
    add_header(slide12, "Level 4B: Ablation Study Level-by-Level Analysis", "Level 4B: Ablation")
    
    add_card(slide12, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2), CARD_BG, AMBER_ACCENT)
    tb_abl1 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf_a1 = tb_abl1.text_frame
    tf_a1.word_wrap = True
    
    p = tf_a1.paragraphs[0]
    p.text = "🧪 Feature Concatenation Ablation"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    p.space_after = Pt(10)
    
    a1_pts = [
        "Single-Scale Top Layer Only (7x7):",
        "  • Test Accuracy: 84.41% (Spatial detail lost in downsampling).",
        "Omit Block 3b Vector (v3):",
        "  • Test Accuracy: 91.20% (-7.37% Drop, fine edge contours lost).",
        "Omit Block 5c Vector (v5):",
        "  • Test Accuracy: 93.80% (-4.77% Drop, landmark shape lost).",
        "Full Multi-Scale Fusion (v3 ⊕ v5 ⊕ v7):",
        "  • Test Accuracy: 98.57% (+14.16% Total Accuracy Gain).",
        "Conclusion: Tapping multi-depth layers is critical to FER."
    ]
    for pt in a1_pts:
        p = tf_a1.add_paragraph()
        p.text = pt
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(6)

    add_card(slide12, Inches(6.833), Inches(1.6), Inches(5.7), Inches(5.2), CARD_BG, VIOLET_ACCENT)
    tb_abl2 = slide12.shapes.add_textbox(Inches(7.033), Inches(1.8), Inches(5.3), Inches(4.8))
    tf_a2 = tb_abl2.text_frame
    tf_a2.word_wrap = True
    
    p = tf_a2.paragraphs[0]
    p.text = "📊 Data Augmentation & Shard Ablation"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = VIOLET_ACCENT
    p.space_after = Pt(10)
    
    a2_pts = [
        "Sharded Dataset Ablation (10-12 Epochs):",
        "  • Standard Augmentation (Flip + Rotation): 36.51% (EffNetB0)",
        "  • No Augmentation on Small Shards: 42.43% (EffNetB0)",
        "Key Finding on Sharded Training:",
        "  • When training only classification heads on small shards, augmentation introduces variance requiring more epochs to overcome.",
        "Full Dataset Convergence:",
        "  • On full 7,529 dataset with end-to-end training, augmentation provides crucial generalization against pose shifts."
    ]
    for pt in a2_pts:
        p = tf_a2.add_paragraph()
        p.text = pt
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(6)

    # ==========================================
    # SLIDE 13: Computational Efficiency & Latency
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    apply_bg(slide13)
    add_header(slide13, "Computational Efficiency & Real-Time Edge Trade-offs", "Efficiency")
    
    eff_cards = [
        ("MobileNetV2", "3.4M", "9.6 MB", "12.1 ms", "Best for Ultra-Low Power Edge Nodes", VIOLET_ACCENT),
        ("YOLOv8 Classifier", "4.2M", "12.8 MB", "8.5 ms", "Fastest Inference (Real-Time Video)", EMERALD_ACCENT),
        ("EfficientNet-B0 (Fused)", "6.8M", "21.4 MB", "14.3 ms", "Champion Accuracy (98.57%) Optimal Balance", CYAN_ACCENT),
        ("ResNet-50", "25.6M", "98.2 MB", "24.5 ms", "Heavy Footprint (High Memory Burden)", AMBER_ACCENT)
    ]
    
    for i, (mname, params, size, lat, text, col) in enumerate(eff_cards):
        left = Inches(0.8 + i * 2.95)
        add_card(slide13, left, Inches(1.6), Inches(2.8), Inches(5.2), CARD_BG, col)
        
        tb = slide13.shapes.add_textbox(left + Inches(0.15), Inches(1.8), Inches(2.5), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = mname
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = f"Params: {params}\nFile Size: {size}\nLatency: {lat}"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(10)
        
        p3 = tf.add_paragraph()
        p3.text = text
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(14)

    # ==========================================
    # SLIDE 14: Conclusion
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    apply_bg(slide14)
    add_header(slide14, "Conclusion & Key Research Accomplishments", "Conclusion")
    
    concls = [
        ("01", "Multi-Scale Fusion Superiority", "Concatenating intermediate activation vectors (Block 3b, Block 5c, Top) solves spatial detail degradation, boosting test accuracy to 98.57%.", CYAN_ACCENT),
        ("02", "Empirical Ecosystem Benchmarking", "Demonstrated that YOLOv8 excels in low-latency environments (8.5ms), while Fused EfficientNet-B0 dominates in overall probabilistic precision (0.2186 log loss).", EMERALD_ACCENT),
        ("03", "Biological Explainability Validated", "Grad-CAM and SHAP attribution confirm that predictions strictly align with Ekman's Facial Action Coding System (eyebrow & mouth contours).", AMBER_ACCENT)
    ]
    
    for i, (num, title, desc, col) in enumerate(concls):
        top = Inches(1.6 + i * 1.8)
        add_card(slide14, Inches(0.8), top, Inches(11.733), Inches(1.6), CARD_BG, col)
        
        add_card(slide14, Inches(1.0), top + Inches(0.2), Inches(0.8), Inches(1.2), col, None)
        tb_num = slide14.shapes.add_textbox(Inches(1.0), top + Inches(0.4), Inches(0.8), Inches(0.8))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = num
        p_num.font.size = Pt(20)
        p_num.font.bold = True
        p_num.font.color.rgb = BG_COLOR
        p_num.alignment = PP_ALIGN.CENTER
        
        tb_desc = slide14.shapes.add_textbox(Inches(2.0), top + Inches(0.15), Inches(10.3), Inches(1.3))
        tf_d = tb_desc.text_frame
        tf_d.word_wrap = True
        
        p = tf_d.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf_d.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 15: Future Scope
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    apply_bg(slide15)
    add_header(slide15, "Future Scope & Strategic Roadmap", "Future Scope")
    
    futures = [
        ("Pillar 1: Hardware Quantization", "INT8 / FP16 TensorRT compilation to enable sub-5ms latency on embedded edge hardware (NVIDIA Jetson / Raspberry Pi 5).", CYAN_ACCENT),
        ("Pillar 2: Multimodal Emotion Fusion", "Combining visual facial expression streams with audio speech acoustics and gait/body language tracking for holistic emotion detection.", EMERALD_ACCENT),
        ("Pillar 3: Spatiotemporal Video Transformers", "Extending static frame classification to 3D-CNNs and Vision Transformers (ViT) to capture continuous micro-expression temporal dynamics.", AMBER_ACCENT)
    ]
    
    for i, (title, desc, col) in enumerate(futures):
        left = Inches(0.8 + i * 3.95)
        add_card(slide15, left, Inches(1.6), Inches(3.75), Inches(5.2), CARD_BG, col)
        
        tb = slide15.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.35), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(14)

    # ==========================================
    # SLIDE 16: References
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    apply_bg(slide16)
    add_header(slide16, "References & Academic Citations", "References")
    
    add_card(slide16, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2), CARD_BG, CYAN_ACCENT)
    tb_ref = slide16.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.8))
    tf_r = tb_ref.text_frame
    tf_r.word_wrap = True
    
    refs = [
        "1. Ekman, P., & Friesen, W. V. (1978). Facial Action Coding System: A Technique for the Measurement of Facial Movement. Consulting Psychologists Press.",
        "2. Tan, M., & Le, Q. V. (2019). EfficientNet: Rethinking Model Scaling for Convolutional Neural Networks. In ICML 2019 (pp. 6105-6114).",
        "3. He, K., Zhang, X., Ren, S., & Sun, J. (2016). Deep Residual Learning for Image Recognition. In IEEE CVPR 2016 (pp. 770-778).",
        "4. Sandler, M., Howard, A., Meng, L., Chen, B., & Adam, H. (2018). MobileNetV2: Inverted Residuals and Linear Bottlenecks. In IEEE CVPR 2018 (pp. 4510-4520).",
        "5. Selvaraju, R. R., Cogswell, M., Das, A., Vedantam, R., Parikh, D., & Batra, D. (2017). Grad-CAM: Visual Explanations from Deep Networks via Gradient-Based Localization. In IEEE ICCV 2017 (pp. 618-626).",
        "6. Jocher, G., et al. (2023). Ultralytics YOLOv8 Architecture & Object Classification Engine. https://github.com/ultralytics/ultralytics"
    ]
    
    for idx, ref in enumerate(refs):
        p = tf_r.add_paragraph() if idx > 0 else tf_r.paragraphs[0]
        p.text = ref
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(10)

    out_file = r"d:\college\DL 4 models\Facial_Emotion_Recognition_Research_Paper_Presentation.pptx"
    prs.save(out_file)
    print(f"Presentation successfully created at: {out_file}")

if __name__ == "__main__":
    create_presentation()
